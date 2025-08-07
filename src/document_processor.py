import requests
import PyPDF2
import docx
import email
import re
from typing import List, Dict
from io import BytesIO
from bs4 import BeautifulSoup
import fitz
import os
from tempfile import NamedTemporaryFile
import logging
import asyncio
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation
import olefile
from src.text_util import clean_escape_characters

pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'

class DocumentProcessor:
    def __init__(self, executor=None):
        self.executor = executor  # Can be None if not using ThreadPoolExecutor
        self.chunk_size = 512
        self.overlap = 50
        self._memory_cache = {}

    async def process_document(self, document_url: str) -> List[Dict]:
        logging.info(f"[DocumentProcessor] Started processing document: {document_url}")

        if document_url in self._memory_cache:
            logging.info("[DocumentProcessor] Returning cached document from memory.")
            return self._memory_cache[document_url]

        loop = asyncio.get_event_loop()

        try:
            # Run blocking requests call in executor
            head = await loop.run_in_executor(self.executor, requests.head, document_url)
            content_type = head.headers.get('content-type', '').lower()
        except Exception:
            logging.warning("[DocumentProcessor] HEAD request failed, skipping content-type check.")
            content_type = None

        logging.info("[DocumentProcessor] Started parsing document")
        if 'pdf' in (content_type or '') or document_url.lower().endswith('.pdf'):
            text = await self._extract_pdf_text_streamed(document_url)
        elif 'spreadsheet' in (content_type or '') or 'excel' in (content_type or '') or document_url.lower().endswith(('.xlsx', '.xls')):
            response = await loop.run_in_executor(self.executor, requests.get, document_url)
            response.raise_for_status()
            text = await loop.run_in_executor(self.executor, self._extract_excel_text, BytesIO(response.content))
        elif 'image' in (content_type or '') or document_url.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
            response = await loop.run_in_executor(self.executor, requests.get, document_url)
            response.raise_for_status()
            text = await loop.run_in_executor(self.executor, self._extract_image_text, BytesIO(response.content))
        elif 'presentation' in (content_type or '') or document_url.lower().endswith(('.pptx', '.ppt')):
            response = await loop.run_in_executor(self.executor, requests.get, document_url)
            response.raise_for_status()
            if document_url.lower().endswith('.pptx'):
                text = await loop.run_in_executor(self.executor, self._extract_pptx_text, BytesIO(response.content))
            else:
                text = await loop.run_in_executor(self.executor, self._extract_ppt_text, BytesIO(response.content))
        else:
            # Run blocking requests call in executor
            response = await loop.run_in_executor(self.executor, requests.get, document_url)
            response.raise_for_status()

            if 'word' in (content_type or '') or document_url.lower().endswith(('.docx', '.doc')):
                text = await loop.run_in_executor(self.executor, self._extract_docx_text, BytesIO(response.content))
            elif 'email' in (content_type or '') or 'message' in (content_type or ''):
                text = await loop.run_in_executor(self.executor, self._extract_email_text, response.content)
            else:
                text = response.text

        logging.info("[DocumentProcessor] Finished parsing document")
        logging.info("[DocumentProcessor] Started text cleaning and chunking")
        # Clean escape characters from extracted text
        text = clean_escape_characters(text)
        chunks = self._create_chunks(text)
        logging.info(f"[DocumentProcessor] Finished chunking. Total chunks created: {len(chunks)}")

        self._memory_cache[document_url] = chunks
        return chunks

    async def _extract_pdf_text_streamed(self, pdf_url: str) -> str:
        logging.info("[DocumentProcessor] Streaming and saving PDF temporarily...")
        loop = asyncio.get_event_loop()
        
        def sync_download_and_parse():
            with requests.get(pdf_url, stream=True) as r:
                r.raise_for_status()
                with NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                    for chunk in r.iter_content(chunk_size=1024 * 1024):
                        temp_file.write(chunk)
                    temp_path = temp_file.name

            logging.info("[DocumentProcessor] Parsing PDF from temp file...")
            text = ""
            try:
                with fitz.open(temp_path) as doc:
                    for page in doc:
                        text += page.get_text()
            finally:
                os.remove(temp_path)
            return text

        return await loop.run_in_executor(self.executor, sync_download_and_parse)

    def _extract_docx_text(self, docx_bytes: BytesIO) -> str:
        try:
            doc = docx.Document(docx_bytes)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            raise Exception(f"Error extracting DOCX text: {str(e)}")

    def _extract_email_text(self, email_content: bytes) -> str:
        try:
            msg = email.message_from_bytes(email_content)
            text = ""

            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        text += part.get_payload(decode=True).decode('utf-8', errors='ignore')
                    elif content_type == "text/html":
                        html_content = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                        soup = BeautifulSoup(html_content, 'html.parser')
                        text += soup.get_text()
            else:
                content_type = msg.get_content_type()
                if content_type == "text/plain":
                    text = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                elif content_type == "text/html":
                    html_content = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text = soup.get_text()

            return text
        except Exception as e:
            raise Exception(f"Error extracting email text: {str(e)}")

    def _extract_excel_text(self, excel_bytes: BytesIO) -> str:
        try:
            df = pd.read_excel(excel_bytes, engine='openpyxl')
            return df.to_string(index=False, header=False)
        except Exception as e:
            raise Exception(f"Error extracting Excel text: {str(e)}")

    def _extract_image_text(self, image_bytes: BytesIO) -> str:
        try:
            image = Image.open(image_bytes)
            text = pytesseract.image_to_string(image)
            print("--- OCR Text ---")
            print(text)
            print("------------------")
            return text
        except Exception as e:
            raise Exception(f"Error extracting image text: {str(e)}")

    def _extract_pptx_text(self, pptx_bytes: BytesIO) -> str:
        try:
            prs = Presentation(pptx_bytes)
            text = []
            for slide in prs.slides:
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                    if shape.shape_type == 13:  # Picture
                        try:
                            image_bytes = BytesIO(shape.image.blob)
                            text.append(self._extract_image_text(image_bytes))
                        except Exception as e:
                            logging.warning(f"Could not extract image from shape: {e}")

                # Check for slide background image
                if slide.background.fill.type == 6: # Picture fill
                    try:
                        image_bytes = BytesIO(slide.background.fill.image.blob)
                        text.append(self._extract_image_text(image_bytes))
                    except Exception as e:
                        logging.warning(f"Could not extract background image: {e}")

            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error extracting PPTX text: {str(e)}")

    def _extract_ppt_text(self, ppt_bytes: BytesIO) -> str:
        text = ""
        try:
            # Try to open as a .pptx file first
            return self._extract_pptx_text(ppt_bytes)
        except Exception:
            # If that fails, handle as a .ppt file
            ppt_bytes.seek(0)
            try:
                ole = olefile.OleFileIO(ppt_bytes)
                
                # Extract text from the main PowerPoint Document stream
                if ole.exists('PowerPoint Document'):
                    ppt_stream = ole.openstream('PowerPoint Document')
                    data = ppt_stream.read()
                    # A simple regex to find long enough ASCII strings
                    for match in re.finditer(b"[\x20-\x7E]{5,}", data):
                        text += match.group(0).decode('ascii', errors='ignore') + "\n"

                # Extract images from picture streams and OCR them
                for stream_name in ole.listdir():
                    if 'Pictures' in stream_name:
                        stream = ole.openstream(stream_name)
                        image_data = stream.read()
                        try:
                            # We don't know the image format, so try to open it with PIL
                            image_bytes = BytesIO(image_data)
                            text += self._extract_image_text(image_bytes)
                        except Exception as e:
                            logging.warning(f"Could not extract image from PPT stream '{stream_name}': {e}")
                return text
            except Exception as e:
                raise Exception(f"Error extracting PPT text: {str(e)}")

    def _create_chunks(self, text: str) -> List[Dict]:
        sentences = re.split(r'[.!?]+', text)
        chunks = []
        current_chunk = ""
        current_length = 0

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            sentence_length = len(sentence.split())

            if current_length + sentence_length > self.chunk_size and current_chunk:
                chunks.append({
                    'text': current_chunk.strip(),
                    'chunk_id': len(chunks),
                    'token_count': current_length
                })

                if len(current_chunk.split()) > self.overlap:
                    overlap_words = current_chunk.split()[-self.overlap:]
                    current_chunk = ' '.join(overlap_words) + ' ' + sentence
                    current_length = len(overlap_words) + sentence_length
                else:
                    current_chunk = sentence
                    current_length = sentence_length
            else:
                current_chunk += ' ' + sentence if current_chunk else sentence
                current_length += sentence_length

        if current_chunk.strip():
            chunks.append({
                'text': current_chunk.strip(),
                'chunk_id': len(chunks),
                'token_count': current_length
            })

        return chunks

    def preprocess_text(self, text: str) -> str:
        # Clean escape characters first
        text = clean_escape_characters(text)
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)]', '', text)
        return text.strip()