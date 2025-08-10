import httpx
import PyPDF2
import docx
import email
import re
from typing import List, Dict, AsyncGenerator
from io import BytesIO
from bs4 import BeautifulSoup
import fitz
import os
from tempfile import NamedTemporaryFile
import logging
import asyncio
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation
import olefile
from src.text_cleaner_utils import clean_escape_characters

# Set the path for the Tesseract OCR executable. This is specific to the environment.
pytesseract.pytesseract.tesseract_cmd = r'C:/Program Files/Tesseract-OCR/tesseract.exe'

class DocumentTextExtractor:
    """
    Processes various document types (PDF, DOCX, Excel, Images, PPTX, Email) from URLs,
    extracts text, and chunks it into smaller, manageable pieces for further processing.
    """
    def __init__(self, executor=None):
        self.executor = executor
        self.chunk_size = 512  # Desired size of text chunks (in words)
        self.overlap = 50      # Number of words to overlap between consecutive chunks
        self.client = httpx.AsyncClient(http2=True, follow_redirects=True) # Async HTTP client

    async def process_document(self, document_url: str) -> AsyncGenerator[Dict, None]:
        """
        Asynchronously processes a document from a given URL.
        Determines document type, extracts text, cleans it, and yields text chunks.
        """
        logging.info(f"Processing document: {document_url}")

        loop = asyncio.get_event_loop()

        # Attempt to get content type from HEAD request, fall back if not available
        content_type = None
        try:
            head = await self.client.head(document_url)
            content_type = head.headers.get('content-type', '').lower()
        except Exception:
            logging.warning("HEAD request failed, skipping content-type check.")

        logging.info("Starting document parsing.")
        text = ""
        # Determine document type based on content-type or file extension and extract text
        if 'pdf' in (content_type or '') or document_url.lower().endswith('.pdf'):
            text = await self._extract_pdf_text_streamed(document_url)
        elif 'spreadsheet' in (content_type or '') or 'excel' in (content_type or '') or document_url.lower().endswith(('.xlsx', '.xls')):
            response = await self.client.get(document_url)
            response.raise_for_status()
            text = await loop.run_in_executor(self.executor, self._extract_excel_text, BytesIO(response.content))
        elif 'image' in (content_type or '') or document_url.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
            response = await self.client.get(document_url)
            response.raise_for_status()
            text = await loop.run_in_executor(self.executor, self._extract_image_text, BytesIO(response.content))
        elif 'presentation' in (content_type or '') or document_url.lower().endswith(('.pptx', '.ppt')):
            response = await self.client.get(document_url)
            response.raise_for_status()
            if document_url.lower().endswith('.pptx'):
                text = await loop.run_in_executor(self.executor, self._extract_pptx_text, BytesIO(response.content))
            else:
                text = await loop.run_in_executor(self.executor, self._extract_ppt_text, BytesIO(response.content))
        else:
            # Default to general text extraction if type is not specifically handled
            response = await self.client.get(document_url)
            response.raise_for_status()

            if 'word' in (content_type or '') or document_url.lower().endswith(('.docx', '.doc')):
                text = await loop.run_in_executor(self.executor, self._extract_docx_text, BytesIO(response.content))
            elif 'email' in (content_type or '') or 'message' in (content_type or ''):
                text = await loop.run_in_executor(self.executor, self._extract_email_text, response.content)
            else:
                text = response.text

        logging.info("Document parsing finished. Starting text cleaning and chunking.")
        text = clean_escape_characters(text)
        # Yield chunks of the extracted and cleaned text
        for chunk in self._create_chunks_generator(text):
            yield chunk

    async def _extract_pdf_text_streamed(self, pdf_url: str) -> str:
        """
        Streams a PDF from a URL, saves it temporarily, extracts text using fitz,
        and then deletes the temporary file.
        """
        logging.info("Streaming and temporarily saving PDF.")
        
        async def async_download_and_parse():
            with NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
                async with self.client.stream("GET", pdf_url) as r:
                    r.raise_for_status()
                    async for chunk in r.aiter_bytes():
                        temp_file.write(chunk)
                temp_path = temp_file.name

            logging.info("Parsing PDF from temporary file.")
            text = ""
            loop = asyncio.get_event_loop()
            try:
                # Run PDF parsing in a separate thread to avoid blocking the event loop
                def parse_pdf():
                    nonlocal text
                    with fitz.open(temp_path) as doc:
                        for page in doc:
                            text += page.get_text()
                await loop.run_in_executor(self.executor, parse_pdf)
            finally:
                # Ensure temporary file is deleted
                await loop.run_in_executor(self.executor, os.remove, temp_path)
            return text

        return await async_download_and_parse()

    def _extract_docx_text(self, docx_bytes: BytesIO) -> str:
        """
        Extracts text from a DOCX file provided as bytes.
        """
        try:
            doc = docx.Document(docx_bytes)
            return "\n".join([p.text for p in doc.paragraphs])
        except Exception as e:
            raise Exception(f"Error extracting DOCX text: {str(e)}")

    def _extract_email_text(self, email_content: bytes) -> str:
        """
        Extracts text from an email file (EML) provided as bytes.
        Handles both plain text and HTML parts.
        """
        try:
            msg = email.message_from_bytes(email_content)
            text = ""

            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        text += part.get_payload(decode=True).decode('utf-8', errors='ignore')
                    elif content_type == "text/html":
                        html_content = part.get_payload(decode=True).decode('utf-8', errors='ignore')
                        soup = BeautifulSoup(html_content, 'html.parser')
                        text += soup.get_text()
            else:
                content_type = msg.get_content_type()
                if content_type == "text/plain":
                    text = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                elif content_type == "text/html":
                    html_content = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text = soup.get_text()

            return text
        except Exception as e:
            raise Exception(f"Error extracting email text: {str(e)}")

    def _extract_excel_text(self, excel_bytes: BytesIO) -> str:
        """
        Extracts text from an Excel file (XLSX/XLS) provided as bytes.
        Converts all sheets to a string representation.
        """
        try:
            df = pd.read_excel(excel_bytes, engine='openpyxl')
            return df.to_string(index=False, header=False)
        except Exception as e:
            raise Exception(f"Error extracting Excel text: {str(e)}")

    def _extract_image_text(self, image_bytes: BytesIO) -> str:
        """
        Extracts text from an image file using Tesseract OCR.
        """
        try:
            image = Image.open(image_bytes)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            raise Exception(f"Error extracting image text: {str(e)}")

    def _extract_pptx_text(self, pptx_bytes: BytesIO) -> str:
        """
        Extracts text from a PPTX (PowerPoint) file provided as bytes.
        Includes text from slides and attempts to extract text from embedded images.
        """
        try:
            prs = Presentation(pptx_bytes)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
                    if shape.shape_type == 13:  # Picture shape type
                        try:
                            image_bytes = BytesIO(shape.image.blob)
                            text.append(self._extract_image_text(image_bytes))
                        except Exception as e:
                            logging.warning(f"Could not extract image from shape: {e}")

                if slide.background.fill.type == 6: # Picture fill type
                    try:
                        image_bytes = BytesIO(slide.background.fill.image.blob)
                        text.append(self._extract_image_text(image_bytes))
                    except Exception as e:
                        logging.warning(f"Could not extract background image: {e}")

            return "\n".join(text)
        except Exception as e:
            raise Exception(f"Error extracting PPTX text: {str(e)}")

    def _extract_ppt_text(self, ppt_bytes: BytesIO) -> str:
        """
        Extracts text from an older PPT (PowerPoint) file provided as bytes.
        Attempts to use PPTX extraction first, then falls back to OLE parsing.
        """
        text = ""
        try:
            # Try to parse as PPTX first, as python-pptx can sometimes handle older formats
            return self._extract_pptx_text(ppt_bytes)
        except Exception:
            ppt_bytes.seek(0)
            try:
                # Fallback to OLE parsing for older PPT formats
                ole = olefile.OleFileIO(ppt_bytes)
                
                if ole.exists('PowerPoint Document'):
                    ppt_stream = ole.openstream('PowerPoint Document')
                    data = ppt_stream.read()
                    # Extract readable strings from the binary data
                    for match in re.finditer(b"[\x20-\x7E]{5,}", data):
                        text += match.group(0).decode('ascii', errors='ignore') + "\n"

                # Attempt to extract text from embedded images within the OLE structure
                for stream_name in ole.listdir():
                    if 'Pictures' in stream_name:
                        stream = ole.openstream(stream_name)
                        image_data = stream.read()
                        try:
                            image_bytes = BytesIO(image_data)
                            text += self._extract_image_text(image_bytes)
                        except Exception as e:
                            logging.warning(f"Could not extract image from PPT stream '{stream_name}': {e}")
                return text
            except Exception as e:
                raise Exception(f"Error extracting PPT text: {str(e)}")

    def _create_chunks_generator(self, text: str):
        """
        Generates text chunks from the input text with specified chunk size and overlap.
        Each chunk is a dictionary containing 'text', 'chunk_id', and 'token_count'.
        """
        sentences = re.split(r'[.!?]+', text) # Split text into sentences
        current_chunk = ""
        current_length = 0
        chunk_id = 0

        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue

            sentence_length = len(sentence.split())

            # If adding the current sentence exceeds chunk size, yield the current chunk
            if current_length + sentence_length > self.chunk_size and current_chunk:
                yield {
                    'text': current_chunk.strip(),
                    'chunk_id': chunk_id,
                    'token_count': current_length
                }
                chunk_id += 1

                # Handle overlap for the next chunk
                if len(current_chunk.split()) > self.overlap:
                    overlap_words = current_chunk.split()[-self.overlap:]
                    current_chunk = ' '.join(overlap_words) + ' ' + sentence
                    current_length = len(overlap_words) + sentence_length
                else:
                    current_chunk = sentence
                    current_length = sentence_length
            else:
                # Add sentence to current chunk
                current_chunk += ' ' + sentence if current_chunk else sentence
                current_length += sentence_length

        # Yield the last chunk if it contains any text
        if current_chunk.strip():
            yield {
                'text': current_chunk.strip(),
                'chunk_id': chunk_id,
                'token_count': current_length
            }

    def preprocess_text(self, text: str) -> str:
        """
        Preprocesses text by cleaning escape characters, normalizing whitespace,
        and removing non-alphanumeric characters (except common punctuation).
        """
        text = clean_escape_characters(text)
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)]', '', text)
        return text.strip()

    async def shutdown(self):
        """
        Closes the asynchronous HTTP client session.
        """
        await self.client.aclose()